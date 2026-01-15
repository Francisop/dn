"""
PPTX Generator Service

Generates professional PowerPoint reports matching NENGIFTOM format.
Based on the working POC from python-poc/generate_slide1_fixed.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from typing import List, Dict, Any
from pathlib import Path
import os

# Brand colors
COLORS = {
    'BROWN': RGBColor(139, 69, 19),
    'ORANGE': RGBColor(255, 192, 0),
    'CREAM': RGBColor(255, 232, 203),
    'WHITE': RGBColor(255, 255, 255),
    'BLACK': RGBColor(0, 0, 0),
    'LIGHTCREAM': RGBColor(255, 244, 231),
}


def format_date_short(date_str: str) -> str:
    """Format date as DD-MM-YYYY"""
    # Already in this format from frontend
    return date_str


def format_date_ordinal(date_input) -> str:
    """Format date to ordinal: '2nd of October, 2024'

    Args:
        date_input: Can be a datetime object or a string in DD-MM-YYYY format
    """
    from datetime import datetime

    # Handle datetime objects directly
    if hasattr(date_input, 'strftime'):
        date_obj = date_input
    else:
        # Parse date string (DD-MM-YYYY)
        try:
            date_obj = datetime.strptime(str(date_input), "%d-%m-%Y")
        except:
            return str(date_input)

    day = date_obj.day
    month = date_obj.strftime("%B")
    year = date_obj.year

    # Ordinal suffix
    if 10 <= day % 100 <= 20:
        suffix = 'th'
    else:
        suffix = {1: 'st', 2: 'nd', 3: 'rd'}.get(day % 10, 'th')

    return f"{day}{suffix} of {month}, {year}"


class PPTXGenerator:
    """PPTX Report Generator"""

    def __init__(self, project_data: Dict[str, Any], incidents: List[Dict[str, Any]],
                 overview_map_path: str = None,
                 satellite_overview_map: List[str] = None,
                 satellite_incident_groups: List[List[Dict[str, Any]]] = None,
                 satellite_thumbnails: List[str] = None,
                 incident_legend_map: str = None):
        """
        Initialize generator with project data

        Args:
            project_data: Project metadata dict
            incidents: List of incident dicts
            overview_map_path: Path to overview map image (optional)
            satellite_overview_map: List of paths to satellite imagery overview maps (optional)
            satellite_incident_groups: List of incident groups (max 3 per group) with metadata (optional)
            satellite_thumbnails: List of paths to satellite thumbnail images (optional)
            incident_legend_map: Path to incident legend image (optional)
        """
        self.project = project_data
        self.incidents = incidents
        self.overview_map_path = overview_map_path
        self.satellite_overview_map = satellite_overview_map or []
        self.satellite_incident_groups = satellite_incident_groups or []
        self.satellite_thumbnails = satellite_thumbnails or []
        self.incident_legend_map = incident_legend_map

        # Logo paths (relative to backend)
        # logo_dir = Path(__file__).parent.parent / "assets" / "logos"
        logo_dir = Path(__file__).resolve()      # current file path
        logo_dir = logo_dir.parents[4] / "assets" / "logos"
        self.logos = str(logo_dir / "logo.png")

    def generate(self, output_path: str) -> str:
        """
        Generate complete PPTX report

        Returns:
            Path to generated PPTX file
        """
        print(f"🎯 Generating PPTX report...")

        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9
        prs.slide_height = Inches(7.5)

        # Set metadata
        prs.core_properties.author = "NENGIFTOM Limited"
        prs.core_properties.title = f"{self.project['projectName']} - Daily Activity Report"
        prs.core_properties.subject = "RPAS Inspection Report"

        # Generate slides
        self._generate_slide1_title(prs)
        self._generate_slide2_executive_summary(prs)
        self._generate_slide3_incident_table(prs)
        self._generate_slide4_map_summary(prs)
        self._generate_slide4_5_summary_table(prs)  # New slide with summary table
        self._generate_slide5_satellite_imagery_map(prs)
        self._generate_analytics_slides(prs)  # NEW: Analytics slides with charts
        # self._generate_slide5_plus_incident_details(prs)  # Now becomes slides 6+

        # Save
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)

        print(f"✅ PPTX saved: {output_path}")
        print(f"   - {len(prs.slides)} slides generated")
        print(f"   - {len(self.incidents)} incidents included")

        return output_path


    def _add_central_watermark(self, slide):
        """Adds the NENGIFTOM logo as a large, semi-transparent, central watermark."""
        if os.path.exists(self.logos):
            watermark_width = Inches(6.0)
            watermark_height = Inches(6.0)

            # Positioned centrally for a 13.333 x 7.5 slide
            left = Inches(4.16)
            top = Inches(1.25)

            pic = slide.shapes.add_picture(self.logos, left, top,
                                          width=watermark_width, height=watermark_height)

            # Make it semi-transparent (10% opacity = 90% transparent)
            # amt value: 10000 = 10% opacity, 50000 = 50% opacity, 100000 = 100% opacity
            try:
                # Access the picture's element
                pic_element = pic._element

                # Define namespaces
                ns = {
                    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
                }

                # Find the blip element (the actual image reference)
                blip = pic_element.find('.//a:blip', namespaces=ns)

                if blip is not None:
                    # Create alpha modulation element using OxmlElement
                    # This sets the transparency level
                    alpha_mod = OxmlElement('a:alphaModFix')
                    alpha_mod.set('amt', '10000')  # 10% opacity = 90% transparent
                    blip.append(alpha_mod)
                    print("✅ Watermark transparency applied (10% opacity)")
                else:
                    print("⚠️ Could not find blip element for transparency")

            except Exception as e:
                print(f"⚠️ Warning: Could not add watermark transparency: {e}")
                import traceback
                traceback.print_exc()
                # Continue without transparency if it fails

    def _add_bullet_to_paragraph(self, paragraph, indent_space_inches=0.35):
        """
        Add bullet point formatting to a paragraph with custom indent spacing.

        When using blank slide layouts, python-pptx doesn't automatically
        add bullet formatting. This manually adds the bullet and controls spacing.

        Args:
            paragraph: The paragraph object to add bullet formatting to
            indent_space_inches: Space between bullet and text in inches (default: 0.35")
        """
        try:
            # Get the paragraph's XML element
            p_element = paragraph._element

            # Get or create paragraph properties (pPr)
            pPr = p_element.get_or_add_pPr()

            # Convert inches to EMUs (English Metric Units)
            # 1 inch = 914400 EMUs
            indent_emu = int(indent_space_inches * 914400)
            margin_emu = indent_emu  # Margin left for bullet position

            # Set marL (margin left) and indent attributes
            # marL: Total left margin for the paragraph
            # indent: Negative value creates hanging indent (space between bullet and text)
            pPr.set('marL', str(margin_emu))
            pPr.set('indent', str(-indent_emu))  # Negative creates the gap

            # Add bullet character
            buChar = OxmlElement('a:buChar')
            buChar.set('char', '•')  # Unicode bullet character
            pPr.append(buChar)

            print(f"✅ Bullet indent set: {indent_space_inches}\" (marL={margin_emu}, indent={-indent_emu})")

        except Exception as e:
            print(f"⚠️ Warning: Could not add bullet formatting: {e}")
            import traceback
            traceback.print_exc()

    def _generate_slide1_title(self, prs):
        """Slide 1: Title Slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_central_watermark(slide)
        # Logos

        slide.shapes.add_picture(self.logos, Inches(0.5), Inches(0.2),
                                width=Inches(0.92), height=Inches(0.92))

        if os.path.exists(self.logos):
            slide.shapes.add_picture(self.logos, Inches(11.913), Inches(0.2),
                                    width=Inches(0.92), height=Inches(0.92))
        # Title
        title_box = slide.shapes.add_textbox(Inches(0), Inches(0.7), Inches(13.333), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = "Daily Activity Report"
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = title_frame.paragraphs[0]
        p.font.size = Pt(60)
        p.font.bold = False
        p.font.name = 'Arial'  # Fallback from Bahnschrift
        p.font.color.rgb = COLORS['BROWN']
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0), Inches(1.9), Inches(13.333), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Project Information"
        subtitle_frame.word_wrap = True
        subtitle_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = subtitle_frame.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.name = 'Arial'
        p.font.color.rgb = COLORS['BROWN']
        p.alignment = PP_ALIGN.CENTER

        # Project details
        # Format date for display
        inspection_date = self.project['inspectionDate']
        date_display = inspection_date.strftime('%d-%m-%Y') if hasattr(inspection_date, 'strftime') else str(inspection_date)

        details_text = f"""Project Name: {self.project['projectName']}
Base Location: {self.project['baseLocation']}
Route Inspected: {self.project['routeInspected']}
Date: {date_display}"""

        details_box = slide.shapes.add_textbox(Inches(1.667), Inches(2.4), Inches(10), Inches(1.5))
        details_frame = details_box.text_frame
        details_frame.text = details_text
        details_frame.word_wrap = True
        details_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        for para in details_frame.paragraphs:
            para.font.size = Pt(21)
            para.font.name = 'Arial'
            para.font.color.rgb = COLORS['BROWN']
            para.alignment = PP_ALIGN.LEFT
            para.line_spacing = 1.5

        # Personnel table (FIXED positioning)
        table_top = Inches(5.4)
        table_height = Inches(1.7)
        table = slide.shapes.add_table(2, 3, Inches(1.667), table_top, Inches(12), table_height).table

        # Column widths
        for col_idx in range(3):
            table.columns[col_idx].width = Inches(3.333)

        # Row heights
        table.rows[0].height = Inches(0.7)
        table.rows[1].height = Inches(1.0)

        # Header row
        headers = ['Prepared by', 'Checked by', 'Approved by']
        for col_idx, header in enumerate(headers):
            cell = table.rows[0].cells[col_idx]
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS['ORANGE']

            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.name = 'Arial'
            p.font.color.rgb = COLORS['WHITE']
            p.alignment = PP_ALIGN.CENTER
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            cell.text_frame.margin_top = Inches(0.05)
            cell.text_frame.margin_bottom = Inches(0.05)

        # Data row
        personnel = [
            self.project['preparedBy'],
            self.project['checkedBy'],
            self.project['approvedBy']
        ]
        for col_idx, person in enumerate(personnel):
            cell = table.rows[1].cells[col_idx]
            cell.text = person
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS['CREAM']

            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(18)
            p.font.name = 'Arial'
            p.font.color.rgb = COLORS['BLACK']
            p.alignment = PP_ALIGN.CENTER
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text_frame.word_wrap = True

            cell.text_frame.margin_top = Inches(0.05)
            cell.text_frame.margin_bottom = Inches(0.05)

    def _generate_slide2_executive_summary(self, prs):
        """Slide 2: Executive Summary"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_central_watermark(slide)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(11), Inches(0.75))
        title_frame = title_box.text_frame
        title_frame.text = "Executive Summary"

        p = title_frame.paragraphs[0]
        p.font.size = Pt(44)
        p.font.bold = False
        p.font.name = 'Arial'
        p.font.color.rgb = COLORS['BROWN']

        # Logo
        if os.path.exists(self.logos):
            slide.shapes.add_picture(self.logos, Inches(11.913), Inches(0.2),
                                    width=Inches(0.92), height=Inches(0.92))

        # Summary text with bold "NENGIFTOM Limited"
        # Use formatted date
        formatted_inspection_date = format_date_ordinal(self.project['inspectionDate'])

        summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(1.2))
        summary_frame = summary_box.text_frame
        summary_frame.word_wrap = True

        # Clear default text and build with runs for mixed formatting
        p = summary_frame.paragraphs[0]

        # First run - "NENGIFTOM Limited" (bold)
        run1 = p.add_run()
        run1.text = "NENGIFTOM Limited"
        run1.font.size = Pt(18)
        run1.font.name = 'Arial'
        run1.font.bold = True
        run1.font.color.rgb = COLORS['BLACK']

        # Second run - rest of the text (not bold)
        run2 = p.add_run()
        run2.text = f" conducted a Remotely Piloted Aircraft System (RPAS) inspection of the {self.project['routeInspected']} Right of Way (RoW) on {formatted_inspection_date}. The total length of the inspected pipeline was {self.project['pipelineLengthKm']}km, and {len(self.incidents)} incident(s) were observed."
        run2.font.size = Pt(18)
        run2.font.name = 'Arial'
        run2.font.color.rgb = COLORS['BLACK']

       # Incident bullets
        if self.incidents:
            # Using Inches(12.333) for width maintains the ~0.5 inch right margin
            # Adjusted top position to 2.7 to clear the summary box (as previously recommended)
            text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(12.333), Inches(4.5)) 
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            
            # NOTE: We can remove text_frame.margin_left = Inches(0.2) 
            # as the indent properties below handle the spacing more precisely.

            for incident in self.incidents:
                p = text_frame.add_paragraph()
                p.text = f"{incident['description']} at coordinates N{incident['latitude']}, E{incident['longitude']}."
                p.level = 0

                # Spacing properties
                p.line_spacing = 1.4        # 1.4x line height (space between lines WITHIN paragraph)
                p.space_after = Pt(12)      # 12 points space AFTER each bullet point

                p.font.size = Pt(18)
                p.font.name = 'Arial'
                p.font.color.rgb = COLORS['BLACK']

                # Add bullet formatting with custom indent spacing (0.35" gap between bullet and text)
                self._add_bullet_to_paragraph(p, indent_space_inches=0.35)

    def _generate_slide3_incident_table(self, prs):
        """Slide 3: Incident Table"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_central_watermark(slide)

        formatted_date = format_date_ordinal(self.project['inspectionDate'])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.2), Inches(0.5), Inches(11.5), Inches(0.75))
        title_frame = title_box.text_frame
        title_frame.text = f"RPAS Activity and Tasking Report for {formatted_date}"
        title_frame.word_wrap = True

        p = title_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.name = 'Arial'
        p.font.color.rgb = COLORS['BROWN']

        # Logo
        if os.path.exists(self.logos):
            slide.shapes.add_picture(self.logos, Inches(11.913), Inches(0.2),
                                    width=Inches(0.92), height=Inches(0.92))

        # Create table
        rows = len(self.incidents) + 1
        cols = 5
        table = slide.shapes.add_table(rows, cols, Inches(0.2), Inches(1.5),
                                       Inches(12.8), Inches(5.5)).table

        # Column widths
        table.columns[0].width = Inches(0.7)  # S/N
        table.columns[1].width = Inches(2.5)    # Incident Points
        table.columns[2].width = Inches(3.8)    # Date/Coordinates
        table.columns[3].width = Inches(0.8)    # N/O
        table.columns[4].width = Inches(5)  # Description

        # Set row heights to 0.73 inches for all rows
        for row in table.rows:
            row.height = Inches(0.73)

        # Header row
        headers = ['S/N', 'Incident Points', 'Date/Coordinates', 'N/O', 'Description']
        for col_idx, header in enumerate(headers):
            cell = table.rows[0].cells[col_idx]
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS['ORANGE']

            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.name = 'Arial'
            p.font.color.rgb = COLORS['WHITE']
            p.alignment = PP_ALIGN.CENTER
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Data rows
        for idx, incident in enumerate(self.incidents):
            row_color = COLORS['CREAM'] if idx % 2 == 0 else COLORS['LIGHTCREAM']
            coords = f"N{incident['latitude']}, E{incident['longitude']}"
            # Format date for display
            date_str = self.project['inspectionDate'].strftime('%d-%m-%Y') if hasattr(self.project['inspectionDate'], 'strftime') else str(self.project['inspectionDate'])
            date_coords = f"{date_str}/{coords}"

            row_data = [
                str(idx + 1),
                incident.get('incidentId', f'INC-{idx+1}'),
                date_coords,
                incident['status'].capitalize(),
                incident['description']
            ]

            for col_idx, data in enumerate(row_data):
                cell = table.rows[idx + 1].cells[col_idx]
                cell.text = data
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_color

                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(16)
                p.font.name = 'Arial'
                p.font.color.rgb = COLORS['BLACK']
                p.alignment = PP_ALIGN.CENTER if col_idx in [0, 3] else PP_ALIGN.LEFT
                cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.text_frame.word_wrap = True

    def _generate_slide4_map_summary(self, prs):
        """Slide 4: Map with Summary Data"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # self._add_central_watermark(slide)

        formatted_date = format_date_ordinal(self.project['inspectionDate'])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.2), Inches(0.2), Inches(8.01), Inches(0.75))
        title_frame = title_box.text_frame
        title_frame.text = f"RPAS Activity and Tasking Report for {formatted_date} Map"
        title_frame.word_wrap = True

        p = title_frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.name = 'Arial'
        p.font.color.rgb = COLORS['BROWN']

        # Logo
        if os.path.exists(self.logos):
            slide.shapes.add_picture(self.logos, Inches(11.913), Inches(0.2),
                                    width=Inches(0.92), height=Inches(0.92))

        # Map (left side)
        if self.overview_map_path and os.path.exists(self.overview_map_path):
            slide.shapes.add_picture(self.overview_map_path, Inches(0.2), Inches(1.3),
                                    width=Inches(8.01), height=Inches(5.95))
        else:
            # Placeholder
            placeholder = slide.shapes.add_textbox(Inches(0.2), Inches(1.5), Inches(8), Inches(5))
            pf = placeholder.text_frame
            pf.text = "[Map will be generated]"
            p = pf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.name = 'Arial'
            p.font.color.rgb = COLORS['BROWN']
            p.alignment = PP_ALIGN.CENTER
            pf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Incident legend (right side) - replaces summary table
        if self.incident_legend_map and os.path.exists(self.incident_legend_map):
            # Legend is 7.2" x 14" (aspect ratio 1:1.944)
            # To maintain aspect ratio: if height = 5.95", then width = 3.06"
            # Center it horizontally in the available space
            legend_width = Inches(3.06)
            legend_height = Inches(5.95)
            legend_x = Inches(8.41) + (Inches(4.92) - legend_width) / 2  # Center horizontally

            slide.shapes.add_picture(self.incident_legend_map, legend_x, Inches(1.3),
                                    width=legend_width, height=legend_height)
        else:
            # Placeholder for legend
            placeholder = slide.shapes.add_textbox(Inches(8.41), Inches(1.5), Inches(4.8), Inches(5))
            pf = placeholder.text_frame
            pf.text = "[Legend will be generated]"
            p = pf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.name = 'Arial'
            p.font.color.rgb = COLORS['BROWN']
            p.alignment = PP_ALIGN.CENTER
            pf.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _generate_slide4_5_summary_table(self, prs):
        """Slide 4.5: Summary Data Table"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_central_watermark(slide)

        # Title
        formatted_date = format_date_ordinal(self.project['inspectionDate'])
        title_text = f"RPAS Activity and Tasking Report for {formatted_date} Summary"

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title_text

        p = title_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.name = 'Arial'
        p.font.color.rgb = COLORS['BROWN']

        # Logo
        if os.path.exists(self.logos):
            slide.shapes.add_picture(self.logos, Inches(11.913), Inches(0.2),
                                    width=Inches(0.92), height=Inches(0.92))

        # Summary table (centered)
        inspection_date = self.project['inspectionDate']
        if hasattr(inspection_date, 'strftime'):
            date_str = inspection_date.strftime('%d-%m-%Y')
        else:
            date_str = str(inspection_date)

        table_data = [
            ["Summary Data", "", COLORS['ORANGE'], True],
            ["Date", date_str, COLORS['ORANGE'], False],
            ["Pipeline RoW Inspected", self.project['routeInspected'], COLORS['ORANGE'], False],
            ["No. of Incident Points Identified", str(len(self.incidents)), COLORS['ORANGE'], False],
            ["Closest Flow stations", self.project['closestFlowStation'], COLORS['ORANGE'], False],
            ["Length of Inspected Pipeline", f"{self.project['pipelineLengthKm']}km", COLORS['ORANGE'], False],
        ]

        # Center the table
        table = slide.shapes.add_table(6, 2, Inches(3), Inches(2),
                                       Inches(7.5), Inches(4.38)).table

        # Set row heights
        for row in table.rows:
            row.height = Inches(0.73)

        for row_idx, (label, value, header_color, is_title) in enumerate(table_data):
            if is_title:
                cell = table.rows[row_idx].cells[0]
                cell.text = label
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                cell.merge(table.rows[row_idx].cells[1])

                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.name = 'Arial'
                p.font.color.rgb = COLORS['WHITE']
                p.alignment = PP_ALIGN.CENTER
                cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            else:
                # Label cell
                label_cell = table.rows[row_idx].cells[0]
                label_cell.text = label
                label_cell.fill.solid()
                label_cell.fill.fore_color.rgb = header_color

                p = label_cell.text_frame.paragraphs[0]
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.name = 'Arial'
                p.font.color.rgb = COLORS['WHITE']
                p.alignment = PP_ALIGN.LEFT
                label_cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                # Value cell
                value_color = COLORS['CREAM'] if (row_idx - 1) % 2 == 0 else COLORS['LIGHTCREAM']
                value_cell = table.rows[row_idx].cells[1]
                value_cell.text = value
                value_cell.fill.solid()
                value_cell.fill.fore_color.rgb = value_color

                p = value_cell.text_frame.paragraphs[0]
                p.font.size = Pt(14)
                p.font.name = 'Arial'
                p.font.color.rgb = COLORS['BLACK']
                p.alignment = PP_ALIGN.LEFT
                value_cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                value_cell.text_frame.word_wrap = True

    def _generate_slide5_satellite_imagery_map(self, prs):
        """Slide 5+: Satellite Imagery Overview with Annotated Incident Images"""

        # Loop through each composite map (map already contains annotated images and connector lines)
        for map_idx, map_path in enumerate(self.satellite_overview_map):
            print(f"🛰️ Generating slide for composite map {map_idx + 1}: {map_path}")

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._add_central_watermark(slide)

            # Title
            formatted_date = format_date_ordinal(self.project['inspectionDate'])
            title_text = f"RPAS Activity and Tasking Report for {formatted_date}"

            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.6))
            title_frame = title_box.text_frame
            title_frame.text = title_text

            p = title_frame.paragraphs[0]
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.name = 'Arial'
            p.font.color.rgb = COLORS['BROWN']

            # Logo
            if os.path.exists(self.logos):
                slide.shapes.add_picture(self.logos, Inches(12.5), Inches(0.3),
                                        width=Inches(0.8), height=Inches(0.8))

            # Add the composite image (already contains map + annotated images + connector lines)
            if map_path and os.path.exists(map_path):
                # Center the composite image on the slide
                slide.shapes.add_picture(map_path, Inches(0.3), Inches(1.2),
                                        width=Inches(12.7), height=Inches(5.9))
                print(f"   ✅ Added composite satellite map with annotations and lines")
            else:
                print(f"   ⚠️ Warning: Composite map not found at {map_path}")

    def _generate_analytics_slides(self, prs):
        """Generate analytics slides with modern charts"""
        import matplotlib
        matplotlib.use("Agg")  # Required for non-GUI environments
        import matplotlib.pyplot as plt
        from collections import Counter
        from .map_generator import categorize_incident

        print("📊 Generating analytics slides...")

        categories = []
        severities = []
        statuses = []

        for incident in self.incidents:
            desc = incident.get('description', '')
            cat = categorize_incident(desc)
            categories.append(cat["name"])

            severities.append(incident.get("severity", "Unknown"))
            statuses.append(incident.get("status", "Unknown"))

        self._create_category_chart_slide(prs, categories)
        self._create_severity_status_slide(prs, severities, statuses)

        print("✅ Analytics slides generated")


    # =====================================================================
    # 🔹 CATEGORY CHART SLIDE (MODERN DONUT + ROUNDED BAR CHART)
    # =====================================================================

    def _create_category_chart_slide(self, prs, categories):
        import matplotlib.pyplot as plt
        from collections import Counter
        import tempfile
        import os
        from .map_generator import categorize_incident

        plt.style.use("fivethirtyeight")  # Modern clean theme

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_central_watermark(slide)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "Incident Analytics: Distribution by Category"
        p = title_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['BROWN']

        # Logo
        if os.path.exists(self.logos):
            slide.shapes.add_picture(self.logos, Inches(12.5), Inches(0.3),
                                    width=Inches(0.8), height=Inches(0.8))

        # Count categories
        category_counts = Counter(categories)
        labels = list(category_counts.keys())
        values = list(category_counts.values())

        # Gather category colors from the categorizer
        category_colors = []
        for name in labels:
            for incident in self.incidents:
                desc = incident.get("description", "")
                cat = categorize_incident(desc)
                if cat["name"] == name:
                    category_colors.append(cat["color"])
                    break

        fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(13, 5.5))

        # ───────────────────────────────
        # MODERN DONUT CHART
        # ───────────────────────────────
        wedges, texts, autotexts = ax1.pie(
            values,
            labels=labels,
            colors=category_colors,
            autopct='%1.1f%%',
            startangle=90,
            pctdistance=0.8,
            wedgeprops=dict(width=0.45, edgecolor='white')
        )
        ax1.set_title("Category Distribution", fontsize=18, weight='bold')

        # White text inside wedges
        for a in autotexts:
            a.set_color("white")
            a.set_weight("bold")
            a.set_fontsize(8)

        # ───────────────────────────────
        # MODERN ROUNDED BAR CHART
        # ───────────────────────────────
        bars = ax2.bar(
            labels,
            values,
            color=category_colors,
            edgecolor="white",
            linewidth=1.5
        )

        # Rounded bars & clean text
        for bar in bars:
            bar.set_alpha(0.92)

        ax2.set_title("Incident Count by Category", fontsize=18, weight="bold")
        ax2.set_ylabel("Count")
        ax2.grid(axis="y", alpha=0.3)

        # Add counts above bars
        for i, val in enumerate(values):
            ax2.text(i, val + 0.1, str(val),
                    ha='center', va='bottom', fontsize=12, weight='bold')

        plt.tight_layout()

        # Save temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        temp_file.close()
        plt.savefig(temp_file.name, dpi=150, bbox_inches="tight", transparent=True)
        plt.close()

        slide.shapes.add_picture(temp_file.name, Inches(0.5), Inches(1.2),
                                width=Inches(12.3))

        os.unlink(temp_file.name)


    # =====================================================================
    # 🔹 SEVERITY + STATUS CHART SLIDE (MODERN BAR CHARTS)
    # =====================================================================

    def _create_severity_status_slide(self, prs, severities, statuses):
        import matplotlib.pyplot as plt
        from collections import Counter
        import tempfile
        import os

        plt.style.use("fivethirtyeight")

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_central_watermark(slide)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "Incident Analytics: Severity & Status"
        p = title_frame.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLORS["BROWN"]

        # Count data
        sev_counts = Counter(severities)
        stat_counts = Counter(statuses)

        fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(13, 5.5))

        # Modern severity bar chart
        ax1.bar(sev_counts.keys(), sev_counts.values(), edgecolor="white", linewidth=1.5)
        ax1.set_title("Severity Distribution", fontsize=18, weight="bold")
        ax1.grid(axis="y", alpha=0.3)

        for i, v in enumerate(sev_counts.values()):
            ax1.text(i, v + 0.1, str(v), ha='center', va='bottom', fontsize=12, weight="bold")

        # Modern status bar chart
        ax2.bar(stat_counts.keys(), stat_counts.values(), edgecolor="white", linewidth=1.5)
        ax2.set_title("Status Distribution", fontsize=18, weight="bold")
        ax2.grid(axis="y", alpha=0.3)

        for i, v in enumerate(stat_counts.values()):
            ax2.text(i, v + 0.1, str(v), ha='center', va='bottom', fontsize=12, weight="bold")

        plt.tight_layout()

        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        temp_file.close()
        plt.savefig(temp_file.name, dpi=150, bbox_inches="tight", transparent=True)
        plt.close()

        slide.shapes.add_picture(temp_file.name, Inches(0.5), Inches(1.2),
                                width=Inches(12.3), height=Inches(5.8))

        os.unlink(temp_file.name)


    def _create_severity_status_slide(self, prs, severities, statuses):
        """Create slide with severity and status distribution"""
        import matplotlib.pyplot as plt
        import tempfile
        from collections import Counter

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_central_watermark(slide)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "Incident Analytics: Severity & Status"

        p = title_frame.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.name = 'Arial'
        p.font.color.rgb = COLORS['BROWN']

        # Logo
        if os.path.exists(self.logos):
            slide.shapes.add_picture(self.logos, Inches(12.5), Inches(0.3),
                                    width=Inches(0.8), height=Inches(0.8))

        # Count data
        severity_counts = Counter(severities)
        status_counts = Counter(statuses)

        # Define colors
        severity_colors = {
            'Critical': '#DC143C',
            'High': '#FF6347',
            'Medium': '#FFA500',
            'Low': '#FFD700',
            'Unknown': '#808080'
        }

        status_colors = {
            'NEW': '#4169E1',
            'IN_PROGRESS': '#FFD700',
            'RESOLVED': '#32CD32',
            'CLOSED': '#808080',
            'Unknown': '#A9A9A9'
        }

        # Create charts
        fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(13, 5.5))

        # Severity bar chart
        sev_colors = [severity_colors.get(s, '#808080') for s in severity_counts.keys()]
        bars1 = ax1.bar(range(len(severity_counts)), list(severity_counts.values()),
                       color=sev_colors, edgecolor='black', linewidth=2)
        ax1.set_xticks(range(len(severity_counts)))
        ax1.set_xticklabels(severity_counts.keys(), fontsize=11, weight='bold')
        ax1.set_ylabel('Count', fontsize=12, weight='bold')
        ax1.set_title('Incidents by Severity', fontsize=16, weight='bold', pad=20)
        ax1.grid(axis='y', alpha=0.3)

        # Add count labels
        for i, count in enumerate(severity_counts.values()):
            ax1.text(i, count + 0.1, str(count), ha='center', va='bottom',
                    fontsize=12, weight='bold')

        # Status donut chart
        stat_colors = [status_colors.get(s, '#808080') for s in status_counts.keys()]
        wedges, texts, autotexts = ax2.pie(
            status_counts.values(),
            labels=status_counts.keys(),
            colors=stat_colors,
            autopct='%1.1f%%',
            startangle=90,
            wedgeprops=dict(width=0.5, edgecolor='white', linewidth=2),
            textprops={'fontsize': 11, 'weight': 'bold'}
        )

        for autotext in autotexts:
            autotext.set_color('white')
            autotext.set_fontsize(13)
            autotext.set_weight('bold')

        ax2.set_title('Incidents by Status', fontsize=16, weight='bold', pad=20)

        plt.tight_layout()

        # Save to temp file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
        temp_file.close()
        plt.savefig(temp_file.name, dpi=150, bbox_inches='tight', transparent=True)
        plt.close()

        # Add to slide
        slide.shapes.add_picture(temp_file.name, Inches(0.5), Inches(1.2),
                                width=Inches(12.3), height=Inches(5.8))

        # Cleanup
        try:
            os.unlink(temp_file.name)
        except:
            pass

    def _generate_slide5_plus_incident_details(self, prs):
        """Slides 5+: Individual Incident Details"""
        for incident in self.incidents:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._add_central_watermark(slide)

            coords = f"N{incident['latitude']}, E{incident['longitude']}"

            # Title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(11.5), Inches(0.75))
            title_frame = title_box.text_frame
            title_frame.text = f"Update: {incident['description']} - {coords}"
            title_frame.word_wrap = True

            p = title_frame.paragraphs[0]
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.name = 'Arial'
            p.font.color.rgb = COLORS['BROWN']

            # Logo
            slide.shapes.add_picture(self.logos, Inches(12.5), Inches(0.3),
                                    width=Inches(0.8), height=Inches(0.8))

            # Annotated photo
            if incident.get('annotatedPhotos') and len(incident['annotatedPhotos']) > 0:
                photo_path = incident['annotatedPhotos'][0]
                if os.path.exists(photo_path):
                    slide.shapes.add_picture(photo_path, Inches(0.5), Inches(1.2),
                                            width=Inches(13), height=Inches(5.5))
                else:
                    # Placeholder
                    ph = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(13), Inches(5.5))
                    phf = ph.text_frame
                    phf.text = f"[Photo: {os.path.basename(photo_path)}]"
                    p = phf.paragraphs[0]
                    p.font.size = Pt(20)
                    p.alignment = PP_ALIGN.CENTER
                    phf.vertical_anchor = MSO_ANCHOR.MIDDLE
            else:
                # Placeholder
                ph = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(13), Inches(5.5))
                phf = ph.text_frame
                phf.text = "[Annotated photo will be added]"
                p = phf.paragraphs[0]
                p.font.size = Pt(24)
                p.font.name = 'Arial'
                p.font.color.rgb = COLORS['BROWN']
                p.alignment = PP_ALIGN.CENTER
                phf.vertical_anchor = MSO_ANCHOR.MIDDLE


async def generate_pptx_report(project_data: Dict[str, Any], incidents: List[Dict[str, Any]],
                               overview_map_path: str, output_path: str,
                               satellite_overview_map: List[str] = None,
                               satellite_incident_groups: List[List[Dict[str, Any]]] = None,
                               satellite_thumbnails: List[str] = None,
                               incident_legend_map: str = None) -> str:
    """
    Generate PPTX report (async wrapper)

    Args:
        project_data: Project metadata
        incidents: List of incidents
        overview_map_path: Path to overview map
        output_path: Output file path
        satellite_overview_map: List of paths to satellite overview maps (optional)
        satellite_incident_groups: List of incident groups with metadata (optional)
        satellite_thumbnails: List of satellite thumbnail paths (optional)
        incident_legend_map: Path to incident legend map (optional)

    Returns:
        Path to generated PPTX
    """
    generator = PPTXGenerator(project_data, incidents, overview_map_path,
                             satellite_overview_map, satellite_incident_groups,
                             satellite_thumbnails, incident_legend_map)
    return generator.generate(output_path)
